import streamlit as st
from PyPDF2 import PdfReader
from PIL import Image
import io
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
import qrcode


# Define a function to extract the first page of a PDF file as an image
def extract_image_from_pdf(pdf_file):
    # Load the PDF file using PyPDF2
    with pdf_file as f:
        pdf = PdfReader(f)
        # Get the first page of the PDF as an image
        page = pdf.pages[0]
        text = page.extract_text()
        
        #get applicant identifiers
        aamc_id = text.split(' ')[2]
        pattern = r"\((\d+)\)"
        aamc_id = re.search(pattern, text).group(1)
        
        full_name = text.split(' ')[0] + " " + text.split(' ')[1]
        
        suffix = "Location:"
        prefix = "Most Recent Medical School:"
        recent_med_school = text.split(suffix)[0].strip()
        med_school = recent_med_school.split(prefix)[1].strip()
    
        # st.write(med_school)
        # st.write(aamc_id)
        # st.write(full_name)
        
        resources = pdf.trailer["/Root"].get_object()["/Pages"].get_object()["/Kids"][0].get_object()["/Resources"].get_object()
        xObject = resources["/XObject"]
        # Find the first image on the page
        
        for obj in xObject:
            if xObject[obj]["/Subtype"] == "/Image":
                # Get the image data
                image_data = xObject[obj]._data
                # Convert the image data to a PIL Image object
                #pil_image = Image.open(io.BytesIO(image_data)
                return image_data, full_name, aamc_id, med_school
            else:
                return None, full_name, aamc_id, med_school

#create Validate URL 
def validate_input(url_field):
    if url_field is not None:
        equals_index = url_field.find("=", url_field.find("=")+1)
        if equals_index != -1:
            # Extract the substring up to the '=' sign
            url_substring = url_field[:equals_index + 1]
            st.success("Valid Google Form URL Form User Input ID")
            return url_substring
        else:
            st.error("No 'valid URL entry ID with =' sign found in the URL given.")

#RESIZED RAW IMAGE 
def resize_image(image_data):
    # Load the image data into a Pillow Image object
    image = Image.open(io.BytesIO(image_data)).convert('RGB')
   
   
    new_width = int(2.5 * 75)
    new_height = int(3.5 * 75)
    #resize the image
    resized_image = image.resize((new_width, new_height), Image.ANTIALIAS)

    # Convert the resized image back to raw image data
    output_buffer = io.BytesIO()
    resized_image.save(output_buffer, format="PNG")
    output_buffer.seek(0)
    resized_image_data = output_buffer.read()
    return resized_image_data
            
#create a QR Code image from entered URL field
def google_url_qrimage(url_substring, full_name):
    google_name_entry = url_substring + full_name

    # Generate the QR code
    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=2,
        border=4
        )
    qr.add_data(google_name_entry)
    qr.make(fit=True)
    # Save the QR code as an image 
    qr_img = qr.make_image(fill_color="black", back_color="white")
    # Convert the image to bytes
    qr_img_bytes = io.BytesIO()
    qr_img.save(qr_img_bytes, format="PNG")
    qr_img_data = qr_img_bytes.getvalue()
    return qr_img_data
    
# Define the Streamlit app
def main():
 

    st.title("PDF Image Extractor and QR Generator to PowerPoint Slides")
    st.write(f"Last update:4/06/23 [Phillip Kim, MD, MPH](https://www.doximity.com/pub/phillip-kim-md-8dccc4e4)")
    st.info('Convert ERAS applicant summary facesheet into PowerPoint slides for custom Rank Meetings and Evaluations')
    pdws_url = 'https://auth.aamc.org/account/#/login?gotoUrl=http:%2F%2Fpdws.aamc.org%2Feras-pdws-web%2F'
    google_url = 'https://www.google.com/forms/about/'
    st.markdown(f"""
    The purpose of this Web App will allow programs to convert applicant facesheet into active user modified post-interview datasheet using QR based system identification into Google Forms.\n
    **ERAS Steps**
    1. login into [PDWS]({pdws_url}) 
    2. Go to **Applications** and select **Current Results**
    3. Use the select box or all to highlight desired applicants
    4. Select **Action to perfrom on slected applicants:** options and Choose **View/Print Application**
    5. Look for in **Other** section **Applicant Summary** which contains applicant photo, AAMC ID, and Medical School information
    6. Select **Print each application to separate PDF**
    7. Type in Print Job Name of choosing then select **Request Print**
    8. Go back to **Bulk Print Request** on output status.  Note: Depending on applicant requests and server demand, ERAS will take a while for this
    9. Save the Zip compiled applicant to your local drive.  Unzip and upload in bulk below.  Note: if you want alphabetical order for PowerPoint slides, you will have to filter by Name in your folder option upon upload\n
    **IMPORTANT: Google Forms**
    1. Create a Google Account if you don't have one. For added security, I would consider 2 Factor Authentication to login as forms generated are stored in Google's server
    2. Head over to [Google Forms]({google_url}) and sign-in
    3. Create **Blank** form
    4. MUST: Change the default *Untitled Question* to *Applicant Name* which will be converted to Short answer format. 
    5. Now feel free to add any other sections that your program will decide on post-interview evaluations.
    6. IMPORTANT: Click on the **3 vertical dots** next to Send Button and select **Get pre-filled link**
    7. In your **Applicant Name** type in 'TEST'
    8. Click **Get link** and to get your **COPY LINK** which is NECESSARY to generate unique QR code for each applicant.  Once a QR code is generated in the PowerPoint slide, this will auto-fill the Google Form when using Mobile Platform
    9. Paste EXACTLY WITHOUT MODIFICATION into **Enter Google Forms URL** below
    """
    )

    # Create a file uploader using Streamlit's file_uploader widget
    pdf_files = st.file_uploader("Upload PDF file(s)", type=["pdf"], accept_multiple_files=True, key='files')

    url_field = st.text_input("Enter Google Forms URL")
    url_substring = validate_input(url_field)

    if url_substring and st.session_state['files']:
    # Process each PDF files
    # create a loop for each upload PDF
    #create presentation file to be appended
        prs = Presentation()
        for index, pdf_file in enumerate(pdf_files):   
            
            # Extract the image from the PDF file
            image_data, full_name, aamc_id, med_school = extract_image_from_pdf(pdf_file)
            # st.write(full_name)
            # st.write(aamc_id)
            
            #resize the image 
            image_data = resize_image(image_data)
            # # Show the image in Streamlit using Streamlit's image widget
            # st.image(image_data, caption="First page of PDF")
            # PASS url_field into function to get image_data
            
            qr_image = google_url_qrimage(url_substring, full_name)
            name_url = url_substring + full_name

            print (url_substring+full_name)
            # Set the properties for the slide numbers
            total_slides = len(pdf_files)
            slide_number_font = 'Arial'
            slide_number_font_size = Pt(12)
            slide_number_color = RGBColor(128, 128, 128)
            slide_number_position = (Inches(8.5), Inches(7))
                
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            top = Inches(2)
            width = Inches(4)
            height = Inches(2)

            # Add applicant image 
            if image_data:
                slide.shapes.add_picture(io.BytesIO(image_data), Inches(1), top)
            # Add QR URL image
            if qr_image:
                slide.shapes.add_picture(io.BytesIO(qr_image), Inches(5), Inches(2))
            
            # Add the applicant name and medical school to the slide
            # Add a text box to the slide
            text_box = slide.shapes.add_textbox(Inches(5), top + Inches(2), width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = f"{full_name}\n{med_school}\nAAMC ID: {aamc_id}"     
            
            # Add a Google forms hyperlink to the full_name variable
            paragraph = text_frame.paragraphs[0]
            run = paragraph.runs[0]
            hyperlink = run.hyperlink
            hyperlink.address = name_url


            # Set the font to Arial and the font size to 24 points for the first line
            font = text_frame.paragraphs[0].font
            font.name = 'Arial'
            font.size = Pt(24)

            # Set the font to Arial and the font size to 18 points for the rest of the text
            for paragraph in text_frame.paragraphs[1:]:
                for run in paragraph.runs: 
                    font = run.font
                    font.name = 'Arial'
                    font.size = Pt(16)
                    
            # Add slide numbers to the slide
            slide_number_text = f'{index+1}/{total_slides}'
            slide_number_box = slide.shapes.add_textbox(*slide_number_position, width=Inches(1), height=Inches(0.2))
            slide_number_frame = slide_number_box.text_frame
            slide_number_frame.text = slide_number_text
            slide_number_frame.paragraphs[0].font.name = slide_number_font
            slide_number_frame.paragraphs[0].font.size = slide_number_font_size
            slide_number_frame.paragraphs[0].font.color.rgb = slide_number_color
            slide_number_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                
        # Save the PowerPoint file
        pptx_file = io.BytesIO()
        prs.save(pptx_file)
        pptx_file.seek(0)
            
        with st.spinner ("Converting into PPT format..."):
            # Download the PowerPoint file
            st.download_button(label="Download PowerPoint", data=pptx_file, file_name="present.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

if __name__ == "__main__":
    main()

